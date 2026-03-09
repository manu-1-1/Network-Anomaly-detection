from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT_BLUE  = RGBColor(0x00, 0xB4, 0xD8)   # cyan-blue
ACCENT_GREEN = RGBColor(0x06, 0xD6, 0xA0)   # mint green
ACCENT_RED   = RGBColor(0xEF, 0x47, 0x6F)   # coral red
ACCENT_YELL  = RGBColor(0xFF, 0xD1, 0x66)   # amber
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE0)
DIM_GRAY     = RGBColor(0x6B, 0x7F, 0x8E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # truly blank layout


# ── helpers ──────────────────────────────────────────────────────
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(1)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_divider(slide, t, color=ACCENT_BLUE, l=0.5, w=12.33):
    rect = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.04))
    rect.fill.solid(); rect.fill.fore_color.rgb = color
    rect.line.fill.background()


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)

# accent bar left
add_rect(slide, 0, 0, 0.18, 7.5, fill=ACCENT_BLUE)

# big title
add_text(slide, "Network Anomaly Detection",
         0.5, 1.6, 12, 1.2, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "using Machine Learning",
         0.5, 2.7, 12, 0.8, size=36, bold=False, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_divider(slide, 3.6, color=ACCENT_GREEN)
add_text(slide, "Real-Time IDS powered by UNSW-NB15 Dataset",
         0.5, 3.8, 12, 0.6, size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, italic=True)
add_text(slide, "Mini Project Presentation  •  2026",
         0.5, 5.8, 12, 0.5, size=16, color=DIM_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Workflow (Mandatory) — matches the image exactly
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, RGBColor(0xFF, 0xFF, 0xFF))   # white background like the image

# Title bar
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x1F, 0x49, 0x7D))
add_text(slide, "Project Workflow  (Mandatory)",
         0.3, 0.12, 12.5, 0.85, size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Main bullet
add_text(slide,
         "•  Problem Definition  →  Data Collection  →  Data Preprocessing",
         0.4, 1.3, 12.5, 0.75,
         size=24, bold=False, color=RGBColor(0x1A, 0x1A, 0x1A), align=PP_ALIGN.LEFT)

# Sub-bullets with arrow
steps = [
    "→  EDA  →  Feature Engineering  →  Model Selection",
    "→  Model Training  →  Model Evaluation",
    "→  Interpretation  →  Conclusion",
]
y = 2.15
for step in steps:
    add_text(slide, f"     –  {step}", 0.8, y, 12, 0.55,
             size=22, color=RGBColor(0x1A, 0x1A, 0x1A), align=PP_ALIGN.LEFT)
    y += 0.6

# bottom accent line
add_rect(slide, 0, 7.3, 13.33, 0.2, fill=RGBColor(0x1F, 0x49, 0x7D))


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — Problem Statement
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "The Problem", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=ACCENT_BLUE)
add_divider(slide, 1.15)

bullets = [
    ("🔴", "Cyber attacks are growing in volume and sophistication"),
    ("🔴", "Traditional rule-based IDS systems miss novel / unknown attacks"),
    ("🟡", "Need: An intelligent system that can detect anomalies automatically"),
    ("🟢", 'Key question: "Can ML models reliably distinguish normal vs. malicious traffic?"'),
]
y = 1.4
for icon, txt in bullets:
    add_text(slide, f"{icon}  {txt}", 0.6, y, 12.1, 0.65, size=20, color=WHITE)
    y += 0.8


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — Dataset
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "Dataset Overview — UNSW-NB15", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=ACCENT_GREEN)
add_divider(slide, 1.15, color=ACCENT_GREEN)

items = [
    "Created by the Australian Centre for Cyber Security (ACCS)",
    "Contains real network traffic + synthetic attack scenarios",
    "49 features per flow  (we use 42: 39 numerical + 3 categorical)",
    "9 attack categories: DoS, Fuzzing, Backdoor, Shellcode, Worms, Reconnaissance…",
    "Binary label:  0 = Normal   |   1 = Attack",
    "Files used:  UNSW_NB15_training-set.csv  /  UNSW_NB15_testing-set.csv",
]
y = 1.45
for item in items:
    add_text(slide, f"  ›  {item}", 0.5, y, 12.3, 0.6, size=19, color=LIGHT_GRAY)
    y += 0.72


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — ML Models Used (table)
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "Models Used", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=ACCENT_YELL)
add_divider(slide, 1.15, color=ACCENT_YELL)

# table headers
cols  = [3.8, 2.5, 2.0, 4.5]
heads = ["Model", "Accuracy", "Speed", "Use Case"]
xs    = [0.25, 4.1, 6.65, 8.7]
y_h   = 1.35

for i, (h, x, cw) in enumerate(zip(heads, xs, cols)):
    add_rect(slide, x, y_h, cw - 0.05, 0.5, fill=ACCENT_BLUE)
    add_text(slide, h, x+0.05, y_h, cw-0.1, 0.5, size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

rows = [
    ("Random Forest",                  "★★★★★  Best",  "~2 s load",  "Production scans"),
    ("Decision Tree",                  "★★★★☆",        "Instant",     "Quick checks"),
    ("Logistic Regression",            "★★★☆☆",        "Instant",     "Baseline"),
    ("Logistic + L1 (feat. sel.)",     "★★★☆☆",        "Instant",     "Feature pruning"),
    ("Ensemble (RF 70% + DT 30%)",     "★★★★★  Robust", "~2 s load", "High-confidence detection"),
]

y_r = y_h + 0.55
row_colors = [RGBColor(0x14, 0x2A, 0x3A), RGBColor(0x1A, 0x31, 0x44)]
for ri, row in enumerate(rows):
    rc = row_colors[ri % 2]
    for ci, (cell, x, cw) in enumerate(zip(row, xs, cols)):
        add_rect(slide, x, y_r, cw - 0.05, 0.52, fill=rc)
        add_text(slide, cell, x+0.05, y_r, cw-0.1, 0.52, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    y_r += 0.55


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "System Architecture", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=ACCENT_BLUE)
add_divider(slide, 1.15)

pipeline = [
    ("Live Network Traffic  (psutil)",         ACCENT_BLUE),
    ("Feature Extraction  (42 UNSW-NB15 features)", ACCENT_GREEN),
    ("ML Pipeline  (sklearn: preprocessing → model)", ACCENT_YELL),
    ("Threat Probability Score  (0 – 100 %)",   ACCENT_RED),
    ("Risk Tier:  LOW / MEDIUM / HIGH / CRITICAL", RGBColor(0xC7, 0x7D, 0xFF)),
    ("CLI Output  (rich terminal UI)",          LIGHT_GRAY),
]

box_w, box_h = 9.5, 0.58
lx = (13.33 - box_w) / 2
y  = 1.35

for i, (label, col) in enumerate(pipeline):
    add_rect(slide, lx, y, box_w, box_h, fill=RGBColor(0x12, 0x28, 0x3C), line=col, line_w=Pt(2))
    add_text(slide, label, lx+0.15, y+0.04, box_w-0.3, box_h-0.1,
             size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        add_text(slide, "↓", lx + box_w/2 - 0.2, y + box_h, 0.4, 0.28,
                 size=18, color=DIM_GRAY, align=PP_ALIGN.CENTER)
    y += box_h + 0.3


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — CLI Tool
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "Real-Time Risk Scorer CLI  (cli_risk_scorer.py)",
         0.4, 0.15, 12.5, 0.8, size=28, bold=True, color=ACCENT_GREEN)
add_divider(slide, 1.15, color=ACCENT_GREEN)

# left: features
feats = [
    "Live capture: reads system-wide network I/O via psutil",
    "Sample mode: scores a random row from UNSW-NB15 test set",
    "Watch mode: continuous monitoring every 10 seconds",
    "Ensemble logic: RF weighted 70 %, DT weighted 30 %",
    "🟢  < 25 %  →  LOW RISK",
    "🟡  25–50 %  →  MEDIUM RISK",
    "🔴  50–75 %  →  HIGH RISK",
    "🚨  > 75 %  →  CRITICAL RISK",
]
y = 1.45
for f in feats:
    add_text(slide, f"  ›  {f}", 0.4, y, 7.0, 0.55, size=17, color=LIGHT_GRAY)
    y += 0.6

# right: commands box
add_rect(slide, 7.8, 1.35, 5.1, 3.2, fill=RGBColor(0x08, 0x14, 0x24), line=ACCENT_GREEN, line_w=Pt(1.5))
cmds = [
    "python cli_risk_scorer.py --model rf",
    "python cli_risk_scorer.py --model ensemble",
    "python cli_risk_scorer.py --detailed",
    "python cli_risk_scorer.py --watch",
    "python cli_risk_scorer.py --sample",
]
yc = 1.5
add_text(slide, "Commands:", 7.95, yc, 4.8, 0.4, size=15, bold=True, color=ACCENT_GREEN)
yc += 0.45
for cmd in cmds:
    add_text(slide, cmd, 7.95, yc, 4.8, 0.45, size=13, color=ACCENT_BLUE)
    yc += 0.48


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — Results & Performance
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "Results & Performance", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=ACCENT_RED)
add_divider(slide, 1.15, color=ACCENT_RED)

results = [
    ("Random Forest", "Highest F1 score on UNSW-NB15 test set"),
    ("Ensemble Model", "Best balance of accuracy + robustness across attack types"),
    ("Attack Detection", "DoS, Fuzzing, Backdoor, Shellcode, Reconnaissance"),
    ("Dynamic Hints", "High packet rate → DoS / Flooding"),
    ("", "Large data transfer → Exfiltration"),
    ("", "FTP/SSH service → Brute Force attempt"),
]
y = 1.45
for label, val in results:
    if label:
        add_text(slide, f"  ✦  {label}:", 0.4, y, 4.0, 0.6, size=19, bold=True, color=ACCENT_YELL)
        add_text(slide, val, 4.5, y, 8.3, 0.6, size=19, color=WHITE)
    else:
        add_text(slide, f"         {val}", 0.4, y, 12.5, 0.55, size=18, color=LIGHT_GRAY)
    y += 0.68


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — Challenges & Limitations
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "Challenges & Limitations", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=ACCENT_YELL)
add_divider(slide, 1.15, color=ACCENT_YELL)

challenges = [
    "UNSW-NB15 is lab traffic — real-world distributions may differ",
    "Live capture uses system-wide aggregates, not per-packet flow data",
    "Some features (jitter, TCP RTT) approximated to 0 without packet capture",
    "Decision Tree probabilities are uncalibrated (capped at 90%)",
    "Large model file: random_forest_ids_model.joblib  ≈  63 MB",
]
y = 1.45
for c in challenges:
    add_text(slide, f"  ⚠  {c}", 0.5, y, 12.3, 0.7, size=20, color=LIGHT_GRAY)
    y += 0.85


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Conclusion & Future Work
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "Conclusion & Future Work", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=ACCENT_GREEN)
add_divider(slide, 1.15, color=ACCENT_GREEN)

done = [
    "Built a fully functional IDS CLI tool trained on a real-world dataset",
    "Supports 3 ML models + ensemble with weighted averaging",
    "Works in live capture mode on any Windows machine via psutil",
]
future = [
    "Integrate Scapy/WinPcap for true per-packet flow capture",
    "Add a web dashboard (Flask/Streamlit) for visual monitoring",
    "Extend to deep learning models (LSTM, Autoencoder) for temporal patterns",
    "Export alerts to SIEM (Security Information & Event Management)",
]

add_text(slide, "✅  Accomplished:", 0.5, 1.35, 6.0, 0.45, size=20, bold=True, color=ACCENT_GREEN)
y = 1.85
for d in done:
    add_text(slide, f"  ✔  {d}", 0.7, y, 11.8, 0.55, size=18, color=WHITE)
    y += 0.6

y += 0.15
add_text(slide, "🚀  Future Work:", 0.5, y, 6.0, 0.45, size=20, bold=True, color=ACCENT_BLUE)
y += 0.5
for f in future:
    add_text(slide, f"  →  {f}", 0.7, y, 11.8, 0.55, size=18, color=LIGHT_GRAY)
    y += 0.58


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — References
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x06, 0x2A, 0x4A))
add_text(slide, "References", 0.4, 0.15, 12, 0.8, size=34, bold=True, color=DIM_GRAY)
add_divider(slide, 1.15, color=DIM_GRAY)

refs = [
    "Moustafa, N. & Slay, J. (2015). UNSW-NB15: A comprehensive data set for network intrusion detection systems. MILCIS.",
    "scikit-learn documentation:  https://scikit-learn.org",
    "psutil documentation:  https://psutil.readthedocs.io",
    "UNSW ACCS Dataset:  https://research.unsw.edu.au/projects/unsw-nb15-dataset",
]
y = 1.5
for r in refs:
    add_text(slide, f"  [{refs.index(r)+1}]  {r}", 0.5, y, 12.3, 0.7, size=18, color=LIGHT_GRAY)
    y += 0.9


# ── Save ────────────────────────────────────────────────────────
out = r"d:\Network-Anomaly-detection\Network_Anomaly_Detection_Presentation.pptx"
prs.save(out)
print(f"✅  Saved: {out}")
